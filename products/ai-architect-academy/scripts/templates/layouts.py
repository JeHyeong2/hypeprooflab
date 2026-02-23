"""Slide layout builder functions using python-pptx."""
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from templates.styles import *


def _set_font(run, size, color=PRIMARY, bold=False, font_name=FONT_FAMILY):
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def _add_textbox(slide, left, top, width, height, text, size, color=PRIMARY,
                 bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size, color, bold)
    return txBox


def _add_rect(slide, left, top, width, height, fill_color=LIGHT_GRAY):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def build_center_title(slide, content):
    """Title slide: centered title + subtitle + optional tagline."""
    # Background accent bar at top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT)

    # Title
    _add_textbox(
        slide, MARGIN_LEFT, Inches(2.5), CONTENT_WIDTH, Inches(1.5),
        content["title"], TITLE_SIZE, PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    # Subtitle
    if content.get("subtitle"):
        _add_textbox(
            slide, MARGIN_LEFT, Inches(4.0), CONTENT_WIDTH, Inches(1.0),
            content["subtitle"], SUBTITLE_SIZE, DARK_GRAY,
            alignment=PP_ALIGN.CENTER
        )
    # Bottom info
    if content.get("bottom"):
        _add_textbox(
            slide, MARGIN_LEFT, Inches(6.2), CONTENT_WIDTH, Inches(0.5),
            content["bottom"], CAPTION_SIZE, MID_GRAY,
            alignment=PP_ALIGN.CENTER
        )


def build_three_column_stats(slide, content):
    """Data slide with 3 big numbers."""
    # Title
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8),
        content["title"], Pt(28), PRIMARY, bold=True
    )
    # Accent underline
    _add_rect(slide, MARGIN_LEFT, Inches(1.5), Inches(2), Inches(0.06), ACCENT)

    cols = content.get("columns", [])
    col_width = Inches(3.2)
    gap = Inches(0.5)
    start_x = MARGIN_LEFT + Inches(0.3)
    y_top = Inches(2.3)

    for i, col in enumerate(cols[:3]):
        x = start_x + i * (col_width + gap)
        # Number
        _add_textbox(
            slide, x, y_top, col_width, Inches(1.2),
            col["stat"], STAT_SIZE, ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Label
        _add_textbox(
            slide, x, y_top + Inches(1.2), col_width, Inches(0.6),
            col["label"], BODY_SIZE, PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Description
        if col.get("desc"):
            _add_textbox(
                slide, x, y_top + Inches(1.8), col_width, Inches(0.8),
                col["desc"], SMALL_SIZE, DARK_GRAY,
                alignment=PP_ALIGN.CENTER
            )

    # Source
    if content.get("source"):
        _add_textbox(
            slide, MARGIN_LEFT, Inches(6.5), CONTENT_WIDTH, Inches(0.4),
            content["source"], CAPTION_SIZE, MID_GRAY
        )


def build_two_column(slide, content):
    """Two-column comparison or explanation."""
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8),
        content["title"], Pt(28), PRIMARY, bold=True
    )
    _add_rect(slide, MARGIN_LEFT, Inches(1.5), Inches(2), Inches(0.06), ACCENT)

    col_width = Inches(5.2)
    gap = Inches(0.8)
    y_start = Inches(2.0)

    for i, col in enumerate([content.get("left", {}), content.get("right", {})]):
        x = MARGIN_LEFT + i * (col_width + gap)
        # Column header
        if col.get("header"):
            _add_textbox(
                slide, x, y_start, col_width, Inches(0.6),
                col["header"], Pt(20), ACCENT if i == 0 else PRIMARY, bold=True
            )
        # Column bullets
        y = y_start + Inches(0.8)
        for bullet in col.get("bullets", []):
            _add_textbox(
                slide, x + Inches(0.2), y, col_width - Inches(0.2), Inches(0.5),
                f"● {bullet}", BODY_SIZE, DARK_GRAY
            )
            y += Inches(0.55)

    if content.get("callout"):
        # Callout box at bottom
        _add_rect(slide, MARGIN_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(1.0), LIGHT_GRAY)
        _add_textbox(
            slide, MARGIN_LEFT + Inches(0.3), Inches(5.95), CONTENT_WIDTH - Inches(0.6), Inches(0.7),
            content["callout"], Pt(14), PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER
        )


def build_bullet_minimal(slide, content):
    """Clean bullet points with optional quote."""
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8),
        content["title"], Pt(28), PRIMARY, bold=True
    )
    _add_rect(slide, MARGIN_LEFT, Inches(1.5), Inches(2), Inches(0.06), ACCENT)

    y = Inches(2.2)
    for bullet in content.get("bullets", []):
        if isinstance(bullet, dict):
            # Bold header + description
            _add_textbox(
                slide, MARGIN_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.4), Inches(0.5),
                f"● {bullet['head']}", Pt(18), PRIMARY, bold=True
            )
            y += Inches(0.5)
            if bullet.get("desc"):
                _add_textbox(
                    slide, MARGIN_LEFT + Inches(0.5), y, CONTENT_WIDTH - Inches(0.7), Inches(0.5),
                    bullet["desc"], BODY_SIZE, DARK_GRAY
                )
                y += Inches(0.55)
        else:
            _add_textbox(
                slide, MARGIN_LEFT + Inches(0.2), y, CONTENT_WIDTH - Inches(0.4), Inches(0.5),
                f"● {bullet}", BODY_SIZE, DARK_GRAY
            )
            y += Inches(0.55)

    if content.get("quote"):
        _add_rect(slide, MARGIN_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(1.0), LIGHT_GRAY)
        _add_textbox(
            slide, MARGIN_LEFT + Inches(0.3), Inches(5.95), CONTENT_WIDTH - Inches(0.6), Inches(0.7),
            f'"{content["quote"]}"', Pt(14), ACCENT, bold=True,
            alignment=PP_ALIGN.CENTER
        )


def build_timeline(slide, content):
    """Timeline/roadmap slide."""
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8),
        content["title"], Pt(28), PRIMARY, bold=True
    )
    _add_rect(slide, MARGIN_LEFT, Inches(1.5), Inches(2), Inches(0.06), ACCENT)

    phases = content.get("phases", [])
    phase_width = Inches(3.3)
    gap = Inches(0.3)
    y_top = Inches(2.3)

    for i, phase in enumerate(phases[:3]):
        x = MARGIN_LEFT + i * (phase_width + gap)
        # Phase box
        _add_rect(slide, x, y_top, phase_width, Inches(3.5), LIGHT_GRAY)
        # Phase name
        _add_textbox(
            slide, x + Inches(0.2), y_top + Inches(0.2), phase_width - Inches(0.4), Inches(0.5),
            phase["name"], Pt(18), ACCENT, bold=True
        )
        # Phase period
        _add_textbox(
            slide, x + Inches(0.2), y_top + Inches(0.7), phase_width - Inches(0.4), Inches(0.4),
            phase.get("period", ""), SMALL_SIZE, MID_GRAY
        )
        # Phase items
        y = y_top + Inches(1.2)
        for item in phase.get("items", []):
            _add_textbox(
                slide, x + Inches(0.3), y, phase_width - Inches(0.6), Inches(0.4),
                f"· {item}", SMALL_SIZE, DARK_GRAY
            )
            y += Inches(0.4)

    if content.get("impact"):
        _add_textbox(
            slide, MARGIN_LEFT, Inches(6.2), CONTENT_WIDTH, Inches(0.5),
            content["impact"], SMALL_SIZE, PRIMARY, bold=True,
            alignment=PP_ALIGN.CENTER
        )


def build_table_slide(slide, content):
    """Table-based slide (e.g., faculty, comparison)."""
    _add_textbox(
        slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.8),
        content["title"], Pt(28), PRIMARY, bold=True
    )
    _add_rect(slide, MARGIN_LEFT, Inches(1.5), Inches(2), Inches(0.06), ACCENT)

    rows = content.get("rows", [])
    headers = content.get("headers", [])
    if not rows:
        return

    num_cols = len(headers) if headers else len(rows[0])
    num_rows = len(rows) + (1 if headers else 0)
    col_w = CONTENT_WIDTH // num_cols

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        MARGIN_LEFT, Inches(2.0),
        CONTENT_WIDTH, Inches(min(num_rows * 0.6, 4.5))
    )
    table = table_shape.table

    # Style header
    if headers:
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, SMALL_SIZE, WHITE, bold=True)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY

    # Fill data
    offset = 1 if headers else 0
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + offset, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, SMALL_SIZE, DARK_GRAY)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


# Layout dispatcher
LAYOUT_BUILDERS = {
    "center-title": build_center_title,
    "three-column-stats": build_three_column_stats,
    "two-column": build_two_column,
    "bullet-minimal": build_bullet_minimal,
    "timeline": build_timeline,
    "table": build_table_slide,
}
