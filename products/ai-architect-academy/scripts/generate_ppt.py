#!/usr/bin/env python3
"""
PPT Generator for Future AI Leader's Academy
Usage: python3 generate_ppt.py [--output path]

Reads slide content definition and generates .pptx using python-pptx.
"""
import sys, os, json
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches
from templates.styles import SLIDE_WIDTH, SLIDE_HEIGHT
from templates.layouts import LAYOUT_BUILDERS

# ============================================================
# Slide content — derived from SPEC.md
# Each slide: { "layout": str, "content": dict }
# ============================================================
SLIDES = [
    # --- Slide 1: Title ---
    {
        "layout": "center-title",
        "content": {
            "title": "Future AI Leader's Academy",
            "subtitle": "코딩하는 아이에서 'AI를 지휘하는 아이'로\n부모와 함께 만드는 4시간의 관점 전환",
            "bottom": "Filamentree (HypeProof AI Lab) × 동아일보"
        }
    },
    # --- Slide 2: Market Hook ---
    {
        "layout": "three-column-stats",
        "content": {
            "title": "시장 기회",
            "columns": [
                {"stat": "430억$", "label": "글로벌 AI 교육", "desc": "2030년 전망 (CAGR 41%)"},
                {"stat": "10배", "label": "한국 코딩교육", "desc": "1,500억→1.5조 (2030)"},
                {"stat": "93%", "label": "학부모 인식", "desc": "AI를 필수 역량으로 인식"},
            ],
            "source": "Sources: Mordor Intelligence, 글로벌인포메이션(GII), 한국경제·입소스"
        }
    },
    # --- Slide 3: Problem Finding ---
    {
        "layout": "bullet-minimal",
        "content": {
            "title": "문제 정의",
            "bullets": [
                {"head": "코딩 기술자 ≠ AI 시대 인재", "desc": "2028년 개발자 75%가 AI 코딩 도구 사용 — 코딩 취업률 90%→60% 급락"},
                {"head": "아이 혼자 하는 교육의 한계", "desc": "부모는 대기실, 가정 내 확장 불가"},
                {"head": "만드는 경험의 부재", "desc": "학교 정보교육 연간 34시간, 이론 위주"},
            ],
            "quote": "코딩 문법을 외우게 할 것인가, AI를 지휘하는 법을 알려줄 것인가?"
        }
    },
    # --- Slide 4: Core Concept ---
    {
        "layout": "two-column",
        "content": {
            "title": "핵심 컨셉: 부모-자녀 Co-Founding",
            "left": {
                "header": "👧 자녀 (Visionary)",
                "bullets": [
                    "무한한 상상력으로 아이디어 발산",
                    '"하늘 나는 고양이를 만들래!"',
                    "AI를 직접 지휘하는 경험",
                ]
            },
            "right": {
                "header": "👨 부모 (Integrator)",
                "bullets": [
                    "현실 감각으로 구조화·구체화",
                    '"날개는 어떤 버튼으로 펴?"',
                    "기획서 정리 지원",
                ]
            },
            "callout": "🤖 AI (Executor) — 지시받은 대로 코드를 작성하는 실행자"
        }
    },
    # --- Slide 5: Curriculum ---
    {
        "layout": "bullet-minimal",
        "content": {
            "title": "커리큘럼: 4시간, 4단계",
            "bullets": [
                {"head": "🔄 Part 1: 관점의 전환 (40분)", "desc": "AI 라이브 데모 + 노동자 vs 건축가 토론"},
                {"head": "🤝 Part 2: 협업과 조율 (50분)", "desc": "샘플 게임 분석 → 기획서 작성 (아이: 상상 / 부모: 정리)"},
                {"head": "🎮 Part 3: 지휘와 소통 (90분)", "desc": "AI 지휘하여 실제 게임 제작, 10초 피드백 루프"},
                {"head": "🏆 Part 4: 성취와 비전 (60분)", "desc": "팀 발표 + 5개 카테고리 시상 (모든 팀 수상)"},
            ]
        }
    },
    # --- Slide 6: Partnership ---
    {
        "layout": "two-column",
        "content": {
            "title": "동아일보 파트너십 구조",
            "left": {
                "header": "동아일보가 하는 것",
                "bullets": [
                    "공동 브랜드 명칭 사용",
                    "동아미디어센터 공간 제공",
                    "어린이동아·동아닷컴 홍보",
                    "구독자·수료생 네트워크 연계",
                ]
            },
            "right": {
                "header": "동아일보가 얻는 것",
                "bullets": [
                    "AI 교육 브랜드 선점",
                    "B2C 신규 매출 (수수료)",
                    "워크숍 기반 콘텐츠 확보",
                    "학부모/학생 데이터",
                ]
            },
            "callout": "추가 인력·예산 투입 없이, 이미 보유한 자산만으로 참여"
        }
    },
    # --- Slide 7: Roadmap ---
    {
        "layout": "timeline",
        "content": {
            "title": "로드맵 & 성과",
            "phases": [
                {
                    "name": "Phase 1: 파일럿",
                    "period": "1~2개월",
                    "items": ["프로그램 검증", "1차 워크숍 진행", "피드백 수집·개선"]
                },
                {
                    "name": "Phase 2: 확산",
                    "period": "3~6개월",
                    "items": ["정기 클래스 운영", "학교 제휴 확보", "월 4회 체계"]
                },
                {
                    "name": "Phase 3: 스케일업",
                    "period": "6~12개월",
                    "items": ["기업 CSR 연계", "강사 양성 프로그램", "온라인 SaaS 확장"]
                }
            ],
            "impact": "NPS 90%+ 목표 | SNS 도달률 회당 3만+ | 손익분기: 월 3회 워크숍"
        }
    },
    # --- Slide 8: Faculty ---
    {
        "layout": "table",
        "content": {
            "title": "Faculty",
            "headers": ["이름", "역할", "소속", "전문분야"],
            "rows": [
                ["이재원 (Jay)", "Product Lead", "Sonatus", "SW Platform, AI Strategy"],
                ["김지웅 (Ryan)", "Tech Lead", "Filamentree", "R&D, Ex-CERN"],
                ["신진용 (JY)", "Content Lead", "Remember", "AI/ML, Agentic AI"],
            ]
        }
    },
    # --- Slide 9: Reference ---
    {
        "layout": "bullet-minimal",
        "content": {
            "title": "실적 & 레퍼런스",
            "bullets": [
                {"head": "워크숍 실적", "desc": "실제 워크숍 진행 완료, 학부모 피드백 수집"},
                {"head": "아이들이 만든 게임 9개", "desc": "4시간 만에 완성한 실제 플레이 가능 게임"},
                {"head": "비즈니스 모델", "desc": "₩200,000/팀, 마진 50%+, 연간 ₩93M 전망"},
            ],
            "quote": "AI가 어렵지 않다는 것을 배웠습니다. 아이와 처음으로 무언가를 함께 완성했습니다."
        }
    },
]


def generate(output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank

    for i, slide_def in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank_layout)
        layout_name = slide_def["layout"]
        builder = LAYOUT_BUILDERS.get(layout_name)
        if not builder:
            print(f"⚠️  Unknown layout '{layout_name}' on slide {i+1}, skipping")
            continue
        builder(slide, slide_def["content"])
        print(f"✅ Slide {i+1}: {layout_name} — {slide_def['content'].get('title', '')}")

    prs.save(output_path)
    print(f"\n🎉 Saved: {output_path}")
    print(f"   {len(SLIDES)} slides generated")


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", "-o",
                        default=os.path.join(os.path.dirname(__file__), "..", "output", "제안서_v2.pptx"))
    args = parser.parse_args()
    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    generate(args.output)
