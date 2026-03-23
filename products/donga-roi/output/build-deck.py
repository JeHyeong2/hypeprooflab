#!/usr/bin/env python3
"""Build deck.pptx from slide-content.json using python-pptx."""

import json, os, io, urllib.request
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ── Design Tokens ────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
CORAL  = RGBColor(0xFF, 0x6B, 0x5A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF5, 0xF7)
MGRAY  = RGBColor(0x8A, 0x8F, 0xA0)
DGRAY  = RGBColor(0x4A, 0x4F, 0x5E)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
LCORAL = RGBColor(0xFF, 0xF0, 0xEE)
FONT   = "Noto Sans KR"

TITLE_SZ = Pt(24)
BODY_SZ  = Pt(14)
SRC_SZ   = Pt(9)
STAT_SZ  = Pt(28)
LABEL_SZ = Pt(11)
SMALL_SZ = Pt(10)

SW = Inches(10)
SH = Inches(5.625)
MX = Inches(0.6)
MY = Inches(0.4)
CW = Inches(8.8)  # 10 - 2*0.6

OUT_DIR = Path(__file__).parent


# ── Image Helpers ────────────────────────────────────────
def fetch_image(url):
    """Download image and return bytes."""
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "pptx-builder"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            return io.BytesIO(resp.read())
    except Exception as e:
        print(f"  WARN: failed to fetch {url}: {e}")
        return None


# ── Text Helpers ─────────────────────────────────────────
def _set_font(run, size=None, color=None, bold=False, italic=False, name=FONT):
    run.font.name = name
    if size: run.font.size = size
    if color: run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_textbox(slide, left, top, width, height, text,
                size=BODY_SZ, color=NAVY, bold=False, italic=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, word_wrap=True):
    """Add a simple single-run text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    _set_font(run, size, color, bold, italic)
    # vertical alignment
    from pptx.oxml.ns import qn
    bodyPr = tf.paragraphs[0]._p.getparent().getparent().find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"
        }.get(valign, "t"))
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape with optional fill."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no line
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


# ── Common Slide Parts ───────────────────────────────────
def add_action_title(slide, text):
    add_textbox(slide, MX, MY, CW, Inches(0.55), text,
                size=TITLE_SZ, color=NAVY, bold=True)


def add_coral_strip(slide):
    add_rect(slide, Inches(0), SH - Inches(0.06), SW, Inches(0.06), CORAL)


def add_slide_number(slide, num):
    add_textbox(slide, SW - Inches(0.8), SH - Inches(0.4), Inches(0.4), Inches(0.3),
                str(num), size=SRC_SZ, color=MGRAY, align=PP_ALIGN.RIGHT)


def add_source_footer(slide, text):
    if not text: return
    add_textbox(slide, MX, SH - Inches(0.4), CW, Inches(0.3),
                "Source: " + text, size=SRC_SZ, color=MGRAY)


def add_callout_bar(slide, left, top, width, height, text, bg=LCORAL, fg=NAVY, size=BODY_SZ):
    """Visible callout bar: light background + dark text."""
    add_rect(slide, left, top, width, height, bg)
    add_textbox(slide, left + Inches(0.3), top, width - Inches(0.6), height,
                text, size=size, color=fg, bold=True, valign=MSO_ANCHOR.MIDDLE)


def std_slide(prs):
    """Create standard white slide with coral strip."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl


# ── Slide Builders ───────────────────────────────────────

def build_s01_cover(prs, s, images):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = NAVY

    # Coral top bar
    add_rect(sl, Inches(0), Inches(0), SW, Inches(0.08), CORAL)

    # Right side: workshop photo (instead of empty decorative rect)
    if images.get("education"):
        pic = sl.shapes.add_picture(images["education"], Inches(7.2), Inches(0.5), Inches(2.5), Inches(4.6))
        # Add semi-transparent overlay for depth
        overlay = add_rect(sl, Inches(7.2), Inches(0.5), Inches(2.5), Inches(4.6), NAVY)
        from pptx.oxml.ns import qn
        from lxml import etree
        solidFill = overlay._element.find(".//" + qn("a:solidFill"))
        if solidFill is not None:
            srgb = solidFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", "40000")  # 40% opacity overlay

    # Logo as text (avoids white bounding box issue)
    add_textbox(sl, MX, Inches(0.5), Inches(3.0), Inches(0.5),
                "HypeProof AI", size=Pt(18), color=CORAL, bold=True)

    # Title (28pt to fit 1 line)
    add_textbox(sl, MX, Inches(1.5), Inches(6.5), Inches(1.0),
                s["action_title"], size=Pt(28), color=WHITE, bold=True)
    # Subtitle
    add_textbox(sl, MX, Inches(2.5), Inches(6.8), Inches(0.5),
                s["subtitle"], size=Pt(20), color=CORAL)
    # Tagline
    add_textbox(sl, MX, Inches(3.2), Inches(6.8), Inches(0.4),
                s["body"]["tagline"], size=BODY_SZ, color=MGRAY, italic=True)
    # Sub-tagline
    add_textbox(sl, MX, Inches(3.6), Inches(6.8), Inches(0.35),
                s["body"]["sub_tagline"], size=Pt(12), color=MGRAY)
    # Partner + Date
    add_textbox(sl, MX, Inches(4.6), CW, Inches(0.3),
                f'{s["body"]["partner"]}  |  {s["body"]["date"]}',
                size=LABEL_SZ, color=MGRAY)


def build_s02_stat_cards(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    cards = s["body"]["cards"]
    cardW = Inches(2.7)
    gap = (CW - len(cards) * cardW) / (len(cards) - 1)
    cardY = Inches(1.2)
    cardH = Inches(3.2)

    for i, c in enumerate(cards):
        cx = MX + i * (cardW + gap)
        add_rect(sl, cx, cardY, cardW, cardH, LGRAY)
        add_rect(sl, cx, cardY, cardW, Inches(0.06), CORAL)

        add_textbox(sl, cx + Inches(0.2), cardY + Inches(0.2), cardW - Inches(0.4), Inches(0.9),
                    c["stat"], size=STAT_SZ, color=NAVY, bold=True)
        add_textbox(sl, cx + Inches(0.2), cardY + Inches(1.15), cardW - Inches(0.4), Inches(0.3),
                    c["label"], size=BODY_SZ, color=CORAL, bold=True)

        # Detail badge
        add_rect(sl, cx + Inches(0.2), cardY + Inches(1.55), cardW - Inches(0.4), Inches(0.35), NAVY)
        add_textbox(sl, cx + Inches(0.2), cardY + Inches(1.55), cardW - Inches(0.4), Inches(0.35),
                    c["detail"], size=SMALL_SZ, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # So-what
        add_textbox(sl, cx + Inches(0.2), cardY + Inches(2.15), cardW - Inches(0.4), Inches(0.85),
                    c["so_what"], size=SMALL_SZ, color=DGRAY, italic=True)

    add_source_footer(sl, s.get("source"))


def build_s03_timeline(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    milestones = s["body"]["milestones"]
    lineY = Inches(2.4)
    startX = MX + Inches(0.5)
    endX = SW - MX - Inches(0.5)
    segW = (endX - startX) / (len(milestones) - 1)

    # Horizontal line
    add_rect(sl, startX, lineY, endX - startX, Inches(0.04), NAVY)

    for i, item in enumerate(milestones):
        cx = startX + int(i * segW)
        isHi = item.get("type") == "highlight"
        dotR = Inches(0.22 if isHi else 0.14)
        dotColor = CORAL if isHi else (MGRAY if item.get("type") == "future" else NAVY)

        add_oval(sl, cx - dotR // 2, lineY - dotR // 2 + Inches(0.02), dotR, dotR, dotColor)

        # Year
        add_textbox(sl, cx - Inches(0.6), lineY - Inches(0.65), Inches(1.2), Inches(0.4),
                    item["year"], size=Pt(18 if isHi else 14), color=CORAL if isHi else NAVY,
                    bold=True, align=PP_ALIGN.CENTER)
        # Event
        add_textbox(sl, cx - Inches(1.0), lineY + Inches(0.3), Inches(2.0), Inches(0.65),
                    item["event"], size=SMALL_SZ, color=NAVY if isHi else DGRAY,
                    bold=isHi, align=PP_ALIGN.CENTER)

        if isHi and item.get("so_what"):
            add_callout_bar(sl, cx - Inches(1.4), lineY + Inches(1.0), Inches(2.8), Inches(0.55),
                            item["so_what"], bg=LCORAL, fg=NAVY, size=SRC_SZ)

    # Bottom callout
    add_rect(sl, MX, Inches(4.3), CW, Inches(0.45), NAVY)
    add_textbox(sl, MX + Inches(0.3), Inches(4.3), CW - Inches(0.6), Inches(0.45),
                s["body"]["callout"], size=Pt(12), color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_source_footer(sl, s.get("source"))


def build_s04_comparison(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    cols = s["body"]["columns"]
    colW = Inches(4.0)
    colGap = Inches(0.8)
    totalW = len(cols) * colW + colGap
    startX = (SW - totalW) // 2
    colY = Inches(1.2)
    colH = Inches(3.0)

    for i, col in enumerate(cols):
        cx = startX + i * (colW + colGap)
        isBlank = col["status"] == "공백"
        bgColor = LCORAL if isBlank else LGRAY
        accentColor = CORAL if isBlank else NAVY

        add_rect(sl, cx, colY, colW, colH, bgColor)
        add_rect(sl, cx, colY, colW, Inches(0.06), accentColor)

        add_textbox(sl, cx + Inches(0.3), colY + Inches(0.2), colW - Inches(0.6), Inches(0.4),
                    col["header"], size=Pt(18), color=NAVY, bold=True)

        add_rect(sl, cx + Inches(0.3), colY + Inches(0.7), Inches(1.2), Inches(0.3), accentColor)
        add_textbox(sl, cx + Inches(0.3), colY + Inches(0.7), Inches(1.2), Inches(0.3),
                    col["status"], size=SMALL_SZ, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Items as bullet list
        itemY = colY + Inches(1.2)
        for j, item in enumerate(col["items"]):
            add_textbox(sl, cx + Inches(0.3), itemY + Inches(j * 0.35), colW - Inches(0.6), Inches(0.35),
                        f"•  {item}", size=BODY_SZ, color=DGRAY)

    # Callout
    add_rect(sl, MX, Inches(4.5), CW, Inches(0.45), NAVY)
    add_textbox(sl, MX + Inches(0.3), Inches(4.5), CW - Inches(0.6), Inches(0.45),
                s["body"]["callout"], size=BODY_SZ, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_source_footer(sl, s.get("source"))


def build_s05_matrix(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    col_headers = s["body"]["col_headers"]
    rows = s["body"]["rows"]
    tableX = MX
    tableY = Inches(1.15)
    nameW = Inches(1.8)
    cellW = Inches(1.2)
    gapColW = Inches(2.2)
    rowH = Inches(0.45)

    # Column headers
    for i, h in enumerate(col_headers):
        x = tableX + nameW + int(i * cellW)
        add_rect(sl, x, tableY, cellW, rowH, NAVY)
        add_textbox(sl, x, tableY, cellW, rowH, h,
                    size=SMALL_SZ, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Gap column header
    gapX = tableX + nameW + int(len(col_headers) * cellW) + Inches(0.1)
    add_rect(sl, gapX, tableY, gapColW, rowH, NAVY)
    add_textbox(sl, gapX, tableY, gapColW, rowH, "차이점",
                size=SMALL_SZ, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    for ri, row in enumerate(rows):
        ry = tableY + int((ri + 1) * rowH)
        isUs = "우리" in row["name"] or "FALA" in row["name"]
        rowBg = LCORAL if isUs else (LGRAY if ri % 2 == 0 else WHITE)

        add_rect(sl, tableX, ry, nameW, rowH, rowBg)
        add_textbox(sl, tableX + Inches(0.05), ry, nameW - Inches(0.1), rowH,
                    row["name"], size=SMALL_SZ, color=CORAL if isUs else NAVY,
                    bold=isUs, valign=MSO_ANCHOR.MIDDLE)

        for vi, v in enumerate(row["cells"]):
            x = tableX + nameW + int(vi * cellW)
            add_rect(sl, x, ry, cellW, rowH, rowBg)
            symbol = "○" if v == "O" else "✕"
            add_textbox(sl, x, ry, cellW, rowH, symbol,
                        size=BODY_SZ, color=GREEN if v == "O" else RED,
                        bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        gapText = row.get("gap") or ("4축 모두 충족" if isUs else "")
        add_rect(sl, gapX, ry, gapColW, rowH, rowBg)
        add_textbox(sl, gapX + Inches(0.05), ry, gapColW - Inches(0.1), rowH,
                    gapText, size=SRC_SZ, color=CORAL if isUs else MGRAY, valign=MSO_ANCHOR.MIDDLE)

    # Callout
    add_callout_bar(sl, MX, Inches(4.0), CW, Inches(0.5),
                    s["body"]["callout"], bg=LCORAL, fg=NAVY, size=BODY_SZ)
    add_source_footer(sl, s.get("source"))


def build_s06_process(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    steps = s["body"]["steps"]
    hasEdu = images.get("education")
    stepsW = Inches(6.8) if hasEdu else CW
    stepW = (stepsW - Inches(0.3) * (len(steps) - 1)) / len(steps)
    stepY = Inches(1.1)
    stepH = Inches(2.5)

    for i, st in enumerate(steps):
        sx = MX + int(i * (stepW + Inches(0.3)))
        add_rect(sl, sx, stepY, stepW, stepH, LGRAY)

        # Step number circle
        circleSize = Inches(0.5)
        cxc = sx + (stepW - circleSize) // 2
        add_oval(sl, cxc, stepY + Inches(0.12), circleSize, circleSize, CORAL)
        add_textbox(sl, cxc, stepY + Inches(0.12), circleSize, circleSize,
                    str(i + 1), size=Pt(16), color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        add_textbox(sl, sx + Inches(0.1), stepY + Inches(0.7), stepW - Inches(0.2), Inches(0.3),
                    st["title"], size=Pt(12), color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(sl, sx + Inches(0.1), stepY + Inches(1.0), stepW - Inches(0.2), Inches(0.22),
                    st["duration"], size=SRC_SZ, color=CORAL, align=PP_ALIGN.CENTER)
        add_textbox(sl, sx + Inches(0.1), stepY + Inches(1.3), stepW - Inches(0.2), Inches(0.45),
                    st["detail"], size=SRC_SZ, color=DGRAY, align=PP_ALIGN.CENTER)

        # Skill badge
        add_rect(sl, sx + Inches(0.1), stepY + Inches(1.9), stepW - Inches(0.2), Inches(0.28), NAVY)
        add_textbox(sl, sx + Inches(0.1), stepY + Inches(1.9), stepW - Inches(0.2), Inches(0.28),
                    st["skill"], size=SRC_SZ, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # Arrow
        if i < len(steps) - 1:
            add_textbox(sl, sx + stepW, stepY + Inches(0.9), Inches(0.3), Inches(0.5),
                        "→", size=Pt(18), color=CORAL, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Education image
    if hasEdu:
        sl.shapes.add_picture(hasEdu, Inches(7.6), Inches(1.1), Inches(2.0), Inches(2.5))

    # Roles
    roles = s["body"]["roles"]
    roleY = Inches(3.85)
    roleW = CW / 3
    for i, (_, desc) in enumerate(roles.items()):
        add_textbox(sl, MX + int(i * roleW), roleY, roleW, Inches(0.35),
                    desc, size=SRC_SZ, color=NAVY, align=PP_ALIGN.CENTER)

    # Tagline
    add_callout_bar(sl, MX, Inches(4.3), CW, Inches(0.4),
                    s["body"]["tagline"], bg=LCORAL, fg=NAVY, size=Pt(12))


def build_s07_channels(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    channels = s["body"]["primary_channels"]
    cardW = Inches(2.7)
    cardGap = (CW - len(channels) * cardW) / (len(channels) - 1)
    cardY = Inches(1.15)
    cardH = Inches(2.95)

    for i, ch in enumerate(channels):
        cx = MX + int(i * (cardW + cardGap))
        add_rect(sl, cx, cardY, cardW, cardH, LGRAY)
        add_rect(sl, cx, cardY, Inches(0.06), cardH, CORAL)

        # Code badge
        add_rect(sl, cx + Inches(0.2), cardY + Inches(0.15), Inches(0.6), Inches(0.3), NAVY)
        add_textbox(sl, cx + Inches(0.2), cardY + Inches(0.15), Inches(0.6), Inches(0.3),
                    ch["code"], size=SRC_SZ, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, cx + Inches(0.9), cardY + Inches(0.15), cardW - Inches(1.1), Inches(0.3),
                    ch["name"], size=BODY_SZ, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, cx + Inches(0.2), cardY + Inches(0.6), cardW - Inches(0.4), Inches(0.3),
                    ch["model"], size=SMALL_SZ, color=DGRAY)

        metrics = [("단가", ch.get("price", "")), ("매출", ch.get("revenue", "")), ("마진", ch.get("margin", ""))]
        for mi, (k, v) in enumerate(metrics):
            my = cardY + Inches(1.05) + int(mi * Inches(0.38))
            add_textbox(sl, cx + Inches(0.2), my, Inches(0.7), Inches(0.28), k, size=SRC_SZ, color=MGRAY)
            add_textbox(sl, cx + Inches(0.9), my, cardW - Inches(1.2), Inches(0.28), v, size=SMALL_SZ, color=NAVY, bold=True)

        add_textbox(sl, cx + Inches(0.2), cardY + Inches(2.3), cardW - Inches(0.4), Inches(0.5),
                    ch["so_what"], size=SRC_SZ, color=CORAL, italic=True)

    # Year 2 expansion bar (2 lines)
    exp = s["body"]["expansion_strip"]
    add_rect(sl, MX, Inches(4.2), CW, Inches(0.7), NAVY)
    add_textbox(sl, MX + Inches(0.3), Inches(4.2), CW - Inches(0.6), Inches(0.32),
                f'{exp["label"]}  →  {exp["year2_target"]}',
                size=LABEL_SZ, color=CORAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(sl, MX + Inches(0.3), Inches(4.52), CW - Inches(0.6), Inches(0.32),
                "  |  ".join(exp["channels"]),
                size=SRC_SZ, color=WHITE, valign=MSO_ANCHOR.MIDDLE)


def build_s08_cases(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    cases = s["body"]["cases"]
    cardW = Inches(2.7)
    gap = (CW - len(cases) * cardW) / (len(cases) - 1)
    cardY = Inches(1.1)
    cardH = Inches(3.3)

    for i, c in enumerate(cases):
        cx = MX + int(i * (cardW + gap))
        add_rect(sl, cx, cardY, cardW, cardH, LGRAY)

        # Icon badge
        add_rect(sl, cx, cardY, cardW, Inches(0.45), NAVY)
        add_textbox(sl, cx + Inches(0.15), cardY, Inches(0.6), Inches(0.45),
                    c["icon"], size=SMALL_SZ, color=CORAL, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, cx + Inches(0.7), cardY, cardW - Inches(0.85), Inches(0.45),
                    c["name"], size=Pt(13), color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

        # Rows
        rows = [
            ("대상", c["target"]),
            ("장소", c["venue"]),
            ("비용", c["cost_bearer"]),
            ("동아 역할", c["donga_role"]),
            ("동아 수익", c["donga_revenue"]),
            ("부가 가치", c["extra_value"]),
        ]
        for ri, (label, val) in enumerate(rows):
            ry = cardY + Inches(0.55) + int(ri * Inches(0.42))
            add_textbox(sl, cx + Inches(0.15), ry, Inches(0.9), Inches(0.2),
                        label, size=Pt(8), color=MGRAY)
            add_textbox(sl, cx + Inches(1.05), ry, cardW - Inches(1.2), Inches(0.38),
                        val, size=SRC_SZ, color=NAVY, bold=(label == "동아 수익"))

    # Callout
    add_callout_bar(sl, MX, Inches(4.55), CW, Inches(0.45),
                    s["body"]["callout"], bg=LCORAL, fg=NAVY, size=Pt(12))
    add_source_footer(sl, s.get("source"))


def build_s09_timetable(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    timeline = s["body"]["timeline"]
    tableX = MX
    tableY = Inches(1.1)
    colWidths = [Inches(1.0), Inches(3.5), Inches(1.8), Inches(1.5)]  # month, action, dept, effort
    rowH = Inches(0.42)

    # Header row
    headers = ["월", "동아 액션", "담당 부서", "공수"]
    ox = 0
    for hi, h in enumerate(headers):
        add_rect(sl, tableX + ox, tableY, colWidths[hi], rowH, NAVY)
        add_textbox(sl, tableX + ox, tableY, colWidths[hi], rowH,
                    h, size=SMALL_SZ, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        ox += colWidths[hi]

    # Data rows
    effort_colors = {"low": GREEN, "medium": RGBColor(0xF5, 0xA6, 0x23)}
    for ri, item in enumerate(timeline):
        ry = tableY + int((ri + 1) * rowH)
        rowBg = LGRAY if ri % 2 == 0 else WHITE
        ox = 0

        # Month
        add_rect(sl, tableX + ox, ry, colWidths[0], rowH, rowBg)
        add_textbox(sl, tableX + ox, ry, colWidths[0], rowH,
                    item["month"], size=SMALL_SZ, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        ox += colWidths[0]

        # Action
        add_rect(sl, tableX + ox, ry, colWidths[1], rowH, rowBg)
        add_textbox(sl, tableX + ox + Inches(0.1), ry, colWidths[1] - Inches(0.2), rowH,
                    item["action"], size=SMALL_SZ, color=DGRAY, valign=MSO_ANCHOR.MIDDLE)
        ox += colWidths[1]

        # Dept
        add_rect(sl, tableX + ox, ry, colWidths[2], rowH, rowBg)
        add_textbox(sl, tableX + ox, ry, colWidths[2], rowH,
                    item["dept"], size=SRC_SZ, color=MGRAY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        ox += colWidths[2]

        # Effort indicator
        add_rect(sl, tableX + ox, ry, colWidths[3], rowH, rowBg)
        eff_color = effort_colors.get(item["effort"], MGRAY)
        eff_label = "가벼움" if item["effort"] == "low" else "보통"
        # Colored dot + label
        dotSize = Inches(0.14)
        dotX = tableX + ox + Inches(0.3)
        dotY = ry + (rowH - dotSize) // 2
        add_oval(sl, dotX, dotY, dotSize, dotSize, eff_color)
        add_textbox(sl, dotX + Inches(0.2), ry, Inches(0.8), rowH,
                    eff_label, size=SRC_SZ, color=eff_color, valign=MSO_ANCHOR.MIDDLE)

    # Callout bar at bottom
    calloutY = tableY + int((len(timeline) + 1) * rowH) + Inches(0.2)
    add_callout_bar(sl, tableX, calloutY, CW, Inches(0.45),
                    s["body"]["callout"], bg=LCORAL, fg=NAVY, size=Pt(11))


def build_s10_waterfall(prs, s, images):  # was s08
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    wf = s["body"]["waterfall"]
    chartX = MX + Inches(0.5)
    chartY = Inches(1.15)
    chartW = Inches(5.2)
    chartH = Inches(3.1)
    maxVal = 14000

    def scaleY(val):
        return chartY + chartH - int((val / maxVal) * chartH)

    # Gridlines
    for v in [0, 3000, 6000, 9000, 12000]:
        gy = scaleY(v)
        add_rect(sl, chartX, gy, chartW, Inches(0.01), RGBColor(0xE0, 0xE0, 0xE0))
        add_textbox(sl, chartX - Inches(0.8), gy - Inches(0.12), Inches(0.7), Inches(0.25),
                    f"{v // 1000}K", size=SRC_SZ, color=MGRAY, align=PP_ALIGN.RIGHT)

    add_textbox(sl, chartX - Inches(0.8), chartY - Inches(0.25), Inches(0.7), Inches(0.2),
                "(만 원)", size=Pt(8), color=MGRAY, align=PP_ALIGN.RIGHT)

    barCount = len(wf)
    barW = int(chartW / (barCount * 1.6))
    barGap = int(barW * 0.6)
    runningCost = 0

    for i, item in enumerate(wf):
        bx = chartX + int(i * (barW + barGap)) + barGap // 2
        colors = {"cost": RED, "subtotal": DGRAY, "revenue": GREEN, "profit": CORAL}
        barColor = colors.get(item["type"], DGRAY)

        if item["type"] == "cost":
            bottom = runningCost
            runningCost += item["value"]
            y1, y2 = scaleY(runningCost), scaleY(bottom)
        elif item["type"] == "subtotal":
            y1, y2 = scaleY(item["value"]), scaleY(0)
        else:
            y1, y2 = scaleY(item["value"]), scaleY(0)

        h = max(y2 - y1, Inches(0.02))
        add_rect(sl, bx, y1, barW, h, barColor)

        add_textbox(sl, bx - Inches(0.2), y1 - Inches(0.28), barW + Inches(0.4), Inches(0.25),
                    f'{item["value"]:,}', size=SMALL_SZ, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(sl, bx - Inches(0.3), chartY + chartH + Inches(0.05), barW + Inches(0.6), Inches(0.3),
                    item["label"], size=SRC_SZ, color=DGRAY, align=PP_ALIGN.CENTER)

    # Scenarios panel
    stX = Inches(6.5)
    stY = Inches(1.2)
    add_rect(sl, stX, stY, Inches(3.2), Inches(0.35), NAVY)
    add_textbox(sl, stX + Inches(0.15), stY, Inches(3.0), Inches(0.35),
                "시나리오 분석", size=Pt(13), color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

    hs = ["시나리오", "매출", "마진", "전제"]
    ws = [Inches(0.8), Inches(0.8), Inches(0.6), Inches(1.0)]
    for hi, h in enumerate(hs):
        ox = sum(ws[:hi])
        add_rect(sl, stX + ox, stY + Inches(0.4), ws[hi], Inches(0.3), LGRAY)
        add_textbox(sl, stX + ox, stY + Inches(0.4), ws[hi], Inches(0.3),
                    h, size=SRC_SZ, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    for ri, sc in enumerate(s["body"]["scenarios"]):
        isTarget = sc["name"] == "목표"
        vals = [sc["name"], sc["revenue"], sc["margin"], sc["premise"]]
        ry = stY + Inches(0.7) + int(ri * Inches(0.4))
        for ci, v in enumerate(vals):
            ox = sum(ws[:ci])
            add_rect(sl, stX + ox, ry, ws[ci], Inches(0.4), WHITE if ri % 2 == 0 else LGRAY)
            add_textbox(sl, stX + ox, ry, ws[ci], Inches(0.4),
                        str(v), size=SRC_SZ, color=CORAL if isTarget else DGRAY,
                        bold=isTarget, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_callout_bar(sl, stX, stY + Inches(2.0), Inches(3.2), Inches(0.55),
                    s["body"]["callout"], bg=LCORAL, fg=NAVY, size=SRC_SZ)
    add_source_footer(sl, s.get("source"))


def build_s09_bep(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    # Native python-pptx line chart
    data = s["body"]["chart_data"]
    chart_data = CategoryChartData()
    chart_data.categories = [d["month"] for d in data]
    chart_data.add_series("누적 매출", [d["cum_revenue"] for d in data])
    chart_data.add_series("누적 비용", [d["cum_cost"] for d in data])

    chart_frame = sl.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.5), Inches(1.1), Inches(6.0), Inches(3.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style series
    s0 = chart.series[0]  # revenue
    s0.format.line.color.rgb = GREEN
    s0.format.line.width = Pt(2.5)
    s0.smooth = False
    s1 = chart.series[1]  # cost
    s1.format.line.color.rgb = RED
    s1.format.line.width = Pt(2.5)
    s1.smooth = False

    # BEP callout
    bep = s["body"]["bep_highlight"]
    bepX = Inches(7.0)
    add_rect(sl, bepX, Inches(1.3), Inches(2.6), Inches(1.3), NAVY)
    add_textbox(sl, bepX + Inches(0.2), Inches(1.4), Inches(2.2), Inches(0.25),
                "BEP 도달", size=Pt(12), color=MGRAY)
    add_textbox(sl, bepX + Inches(0.2), Inches(1.7), Inches(2.2), Inches(0.45),
                bep["month"], size=Pt(22), color=CORAL, bold=True)
    add_textbox(sl, bepX + Inches(0.2), Inches(2.15), Inches(2.2), Inches(0.35),
                bep["message"], size=SMALL_SZ, color=WHITE)

    # Premises
    add_textbox(sl, bepX, Inches(2.8), Inches(2.6), Inches(0.3),
                "전제 조건", size=Pt(12), color=NAVY, bold=True)
    for pi, premise in enumerate(s["body"]["premises"]):
        add_textbox(sl, bepX + Inches(0.15), Inches(3.1) + int(pi * Inches(0.28)), Inches(2.45), Inches(0.28),
                    f"•  {premise}", size=SRC_SZ, color=DGRAY)

    # So-what
    add_callout_bar(sl, bepX, Inches(4.1), Inches(2.6), Inches(0.5),
                    s["body"]["so_what"], bg=LCORAL, fg=NAVY, size=SRC_SZ)
    add_source_footer(sl, s.get("source"))


def build_s10_give_get(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    left_data, right_data = s["body"]["left"], s["body"]["right"]
    colW = Inches(4.2)
    colGap = Inches(0.4)
    colY = Inches(1.05)
    colH = Inches(3.4)

    # LEFT: what Donga provides
    lx = MX
    add_rect(sl, lx, colY, colW, colH, LGRAY)
    add_rect(sl, lx, colY, colW, Inches(0.5), NAVY)
    add_textbox(sl, lx + Inches(0.2), colY, Inches(2.5), Inches(0.5),
                left_data["header"], size=BODY_SZ, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(sl, lx + Inches(2.5), colY, Inches(1.5), Inches(0.5),
                left_data["sub_header"], size=SMALL_SZ, color=CORAL, valign=MSO_ANCHOR.MIDDLE)

    for i, row in enumerate(left_data["rows"]):
        ry = colY + Inches(0.6) + int(i * Inches(0.55))
        add_textbox(sl, lx + Inches(0.2), ry, Inches(2.2), Inches(0.25),
                    row["item"], size=SMALL_SZ, color=NAVY, bold=True)
        add_textbox(sl, lx + Inches(0.2), ry + Inches(0.22), Inches(2.2), Inches(0.2),
                    "실비용: " + row["real_cost"], size=SRC_SZ, color=GREEN)
        add_textbox(sl, lx + Inches(2.5), ry, Inches(1.5), Inches(0.45),
                    "시가: " + row["value"], size=SRC_SZ, color=MGRAY, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)

    # RIGHT: what Donga gets
    rx = lx + colW + colGap
    add_rect(sl, rx, colY, colW, colH, RGBColor(0xFF, 0xF8, 0xF6))
    add_rect(sl, rx, colY, colW, Inches(0.5), CORAL)
    add_textbox(sl, rx + Inches(0.2), colY, Inches(2.5), Inches(0.5),
                right_data["header"], size=BODY_SZ, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(sl, rx + Inches(2.5), colY, Inches(1.5), Inches(0.5),
                right_data["sub_header"], size=SMALL_SZ, color=NAVY, valign=MSO_ANCHOR.MIDDLE)

    for i, row in enumerate(right_data["rows"]):
        ry = colY + Inches(0.6) + int(i * Inches(0.55))
        add_textbox(sl, rx + Inches(0.2), ry, Inches(1.6), Inches(0.25),
                    row["item"], size=SMALL_SZ, color=NAVY, bold=True)
        add_textbox(sl, rx + Inches(0.2), ry + Inches(0.22), Inches(2.5), Inches(0.2),
                    row["so_what"], size=SRC_SZ, color=MGRAY)
        add_textbox(sl, rx + Inches(2.8), ry, Inches(1.2), Inches(0.22),
                    "Y1: " + row["y1"], size=SRC_SZ, color=CORAL, align=PP_ALIGN.RIGHT)
        add_textbox(sl, rx + Inches(2.8), ry + Inches(0.22), Inches(1.2), Inches(0.2),
                    "Y2: " + row["y2"], size=SRC_SZ, color=NAVY, align=PP_ALIGN.RIGHT)

    # Bottom callout
    add_rect(sl, MX, Inches(4.6), CW, Inches(0.4), NAVY)
    add_textbox(sl, MX + Inches(0.3), Inches(4.6), CW - Inches(0.6), Inches(0.4),
                s["body"]["callout"], size=Pt(12), color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_source_footer(sl, s.get("source"))


def build_s11_team(prs, s, images):
    sl = std_slide(prs)
    add_action_title(sl, s["action_title"])
    add_coral_strip(sl)
    add_slide_number(sl, s["_slide_num"])

    members = s["body"]["members"]
    cols = 3
    cardW = Inches(2.7)
    cardH = Inches(1.35)
    gapX = (CW - cols * cardW) / (cols - 1)
    gapY = Inches(0.15)
    startY = Inches(1.05)

    team_images = images.get("team", [])

    for i, m in enumerate(members):
        col = i % cols
        row = i // cols
        cx = MX + int(col * (cardW + gapX))
        cy = startY + int(row * (cardH + gapY))

        add_rect(sl, cx, cy, cardW, cardH, LGRAY)

        # Photo or text-only layout
        photoSize = Inches(0.65)
        imgRef = m.get("image")
        teamIdx = -1
        if imgRef and imgRef.startswith("team["):
            try: teamIdx = int(imgRef.split("[")[1].split("]")[0])
            except: pass

        hasPhoto = teamIdx >= 0 and teamIdx < len(team_images) and team_images[teamIdx]
        textX = cx + Inches(0.9) if hasPhoto else cx + Inches(0.2)
        textW = cardW - Inches(1.05) if hasPhoto else cardW - Inches(0.4)

        if hasPhoto:
            sl.shapes.add_picture(team_images[teamIdx], cx + Inches(0.12), cy + Inches(0.15), photoSize, photoSize)
        else:
            # No photo: use colored accent bar on left instead of initial circle
            add_rect(sl, cx, cy, Inches(0.08), cardH, CORAL)

        add_textbox(sl, textX, cy + Inches(0.1), textW, Inches(0.28),
                    m["name"], size=Pt(12), color=NAVY, bold=True)
        add_textbox(sl, textX, cy + Inches(0.38), textW, Inches(0.22),
                    m["role"], size=SMALL_SZ, color=CORAL)
        add_textbox(sl, textX, cy + Inches(0.6), textW, Inches(0.2),
                    m["org"], size=SRC_SZ, color=DGRAY)
        add_textbox(sl, cx + Inches(0.12), cy + Inches(0.88), cardW - Inches(0.24), Inches(0.4),
                    m["credential"], size=Pt(8), color=MGRAY)

    # Proof bar
    pp = s["body"]["proof_points"]
    ppY = Inches(3.85)
    add_rect(sl, MX, ppY, CW, Inches(0.8), NAVY)
    add_textbox(sl, MX + Inches(0.3), ppY + Inches(0.05), CW - Inches(0.6), Inches(0.25),
                pp["label"], size=LABEL_SZ, color=CORAL, bold=True)
    for pi, item in enumerate(pp["items"]):
        add_textbox(sl, MX + Inches(0.3), ppY + Inches(0.3) + int(pi * Inches(0.16)), CW - Inches(0.6), Inches(0.16),
                    f"•  {item}", size=SRC_SZ, color=WHITE)


def build_s12_roadmap(prs, s, images):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = NAVY

    add_rect(sl, Inches(0), Inches(0), SW, Inches(0.06), CORAL)

    add_textbox(sl, MX, Inches(0.25), CW, Inches(0.55),
                s["action_title"], size=TITLE_SZ, color=WHITE, bold=True)

    # Phases
    phases = s["body"]["roadmap"]
    phW = Inches(2.1)
    phGap = Inches(0.1)
    phY = Inches(0.95)
    phH = Inches(1.55)

    for i, ph in enumerate(phases):
        px = MX + int(i * (phW + phGap))
        isFirst = ph.get("highlight", False)
        phColor = CORAL if isFirst else RGBColor(0x2A, 0x3F, 0x6A)

        r = add_rect(sl, px, phY, phW, phH, phColor)
        if not isFirst:
            # Set transparency via XML
            from pptx.oxml.ns import qn
            from lxml import etree
            solidFill = r._element.find(".//" + qn("a:solidFill"))
            if solidFill is not None:
                srgb = solidFill.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha = etree.SubElement(srgb, qn("a:alpha"))
                    alpha.set("val", "50000")

        add_textbox(sl, px + Inches(0.1), phY + Inches(0.08), phW - Inches(0.2), Inches(0.22),
                    ph["phase"], size=SRC_SZ, color=WHITE if isFirst else MGRAY)
        add_textbox(sl, px + Inches(0.1), phY + Inches(0.3), phW - Inches(0.2), Inches(0.28),
                    ph["title"], size=Pt(12), color=WHITE, bold=True)
        add_textbox(sl, px + Inches(0.1), phY + Inches(0.58), phW - Inches(0.2), Inches(0.22),
                    ph["period"], size=SRC_SZ, color=CORAL)
        add_textbox(sl, px + Inches(0.1), phY + Inches(0.78), phW - Inches(0.2), Inches(0.2),
                    ph.get("cost") or "", size=Pt(8), color=MGRAY)
        add_textbox(sl, px + Inches(0.1), phY + Inches(1.0), phW - Inches(0.2), Inches(0.45),
                    ph["key_action"], size=Pt(8), color=RGBColor(0xB0, 0xB8, 0xCC))

    # Asks section
    askY = Inches(2.7)
    add_textbox(sl, MX, askY, CW, Inches(0.35),
                "요청 사항 (ASK)", size=BODY_SZ, color=CORAL, bold=True)

    asks = s["body"]["asks"]
    askW = CW / len(asks)
    for i, a in enumerate(asks):
        ax = MX + int(i * askW)
        add_rect(sl, ax + Inches(0.05), askY + Inches(0.4), askW - Inches(0.1), Inches(1.1),
                 RGBColor(0x2A, 0x3F, 0x6A))

        circleSize = Inches(0.32)
        add_oval(sl, ax + Inches(0.15), askY + Inches(0.5), circleSize, circleSize, CORAL)
        add_textbox(sl, ax + Inches(0.15), askY + Inches(0.5), circleSize, circleSize,
                    str(i + 1), size=Pt(12), color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        add_textbox(sl, ax + Inches(0.55), askY + Inches(0.48), askW - Inches(0.75), Inches(0.35),
                    a["request"], size=SMALL_SZ, color=WHITE, bold=True)
        add_textbox(sl, ax + Inches(0.55), askY + Inches(0.83), askW - Inches(0.75), Inches(0.35),
                    f'{a["timing"]} / {a["detail"]}', size=Pt(8), color=MGRAY)

    # Callout
    add_rect(sl, MX, Inches(4.1), CW, Inches(0.4), CORAL)
    add_textbox(sl, MX + Inches(0.3), Inches(4.1), CW - Inches(0.6), Inches(0.4),
                s["body"]["callout"], size=Pt(12), color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)

    # Logo as text (avoids white bounding box)
    add_textbox(sl, SW - Inches(2.5), SH - Inches(0.55), Inches(2.0), Inches(0.4),
                "HypeProof AI", size=Pt(14), color=CORAL, bold=True, align=PP_ALIGN.RIGHT)

    # Footer
    add_textbox(sl, MX, SH - Inches(0.4), CW - Inches(2.0), Inches(0.3),
                "Filamentree × 동아일보  |  Future AI Leader's Academy",
                size=SRC_SZ, color=MGRAY)


# ── Main ─────────────────────────────────────────────────
def main():
    content = json.loads((OUT_DIR / "slide-content.json").read_text())
    slides = content["slides"]

    # Download images
    print("Downloading images...")
    img = {}
    img["logo"] = fetch_image("https://hypeproof-ai.xyz/logos/logo-h-light-lg.png")

    team_imgs = []
    for url in ["https://hypeproof-ai.xyz/members/jay.png",
                "https://hypeproof-ai.xyz/members/ryan.png",
                "https://hypeproof-ai.xyz/members/jy.png"]:
        team_imgs.append(fetch_image(url))
    img["team"] = team_imgs

    edu = fetch_image("https://hypeproof-ai.xyz/workshop/kid-coding.png")
    img["education"] = edu

    # Build presentation
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    builders = [
        build_s01_cover, build_s02_stat_cards, build_s03_timeline,
        build_s04_comparison, build_s05_matrix, build_s06_process,
        build_s07_channels, build_s08_cases, build_s09_timetable,
        build_s10_waterfall, build_s09_bep, build_s10_give_get,
        build_s11_team, build_s12_roadmap,
    ]

    for i, builder in enumerate(builders):
        print(f"  Building S{i+1:02d}...")
        # Inject slide number into content data for builders that use it
        slides[i]["_slide_num"] = i + 1
        builder(prs, slides[i], img)

    out_path = OUT_DIR / "deck.pptx"
    prs.save(str(out_path))
    print(f"Done: {out_path} ({out_path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
